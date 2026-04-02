from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDES_DIR = os.path.dirname(os.path.abspath(__file__))

# Colors
WARM_BG = RGBColor(0xFF, 0xFC, 0xF7)
WARM_BG2 = RGBColor(0xF5, 0xED, 0xE0)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BROWN = RGBColor(0x6B, 0x5A, 0x4E)
ACCENT = RGBColor(0xC0, 0x78, 0x5A)
LIGHT_ACCENT = RGBColor(0xD4, 0xA5, 0x74)
MUTED = RGBColor(0xB5, 0xA8, 0x98)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_BOX = RGBColor(0xE8, 0xDC, 0xC8)
RED_ACCENT = RGBColor(0xE8, 0x78, 0x5A)
BLUE_ACCENT = RGBColor(0x5B, 0x8F, 0xD4)
TAG_BG = RGBColor(0xF0, 0xD6, 0xD6)
TAG_FG = RGBColor(0x8B, 0x4A, 0x4A)

def set_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = WARM_BG
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[1].color.rgb = WARM_BG2
    fill.gradient_stops[1].position = 1.0

def add_text(slide, left, top, width, height, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box

def add_rich_bullets(slide, left, top, width, height, items, size=14, color=DARK, spacing=8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
        # Add bullet character as run then text
        run = p.add_run()
        run.text = "\u2022  "
        run.font.size = Pt(size)
        run.font.color.rgb = LIGHT_ACCENT
        run.font.name = 'Calibri'
        run.font.bold = True
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(size)
        run2.font.color.rgb = color
        run2.font.name = 'Calibri'
    return box

def add_rounded_box(slide, left, top, width, height, fill_color=WHITE, border_color=None, border_width=1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_callout_box(slide, left, top, width, text, size=13):
    height = max(0.6, len(text) / 80 * 0.4 + 0.4)
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.06), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_ACCENT
    bar.line.fill.background()
    # Box
    shape = add_rounded_box(slide, left + 0.08, top, width - 0.08, height, WHITE)
    shape.adjustments[0] = 0.03
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.italic = True
    p.font.color.rgb = BROWN
    p.font.name = 'Calibri'
    return height

def add_tag(slide, left, top, text, bg=CREAM_BOX, fg=BROWN, width=1.5):
    shape = add_rounded_box(slide, left, top, width, 0.32, bg)
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Calibri'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_flow_box(slide, left, top, width, height, line1, line2, size=12):
    shape = add_rounded_box(slide, left, top, width, height, WHITE, CREAM_BOX)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = line1
    p.font.size = Pt(size)
    p.font.color.rgb = DARK
    p.font.name = 'Calibri'
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if line2:
        p2 = tf.add_paragraph()
        p2.text = line2
        p2.font.size = Pt(size - 1)
        p2.font.color.rgb = BROWN
        p2.font.name = 'Calibri'
        p2.alignment = PP_ALIGN.CENTER

def add_arrow(slide, left, top, size=22):
    add_text(slide, left, top, 0.4, 0.5, "\u2192", size=size, color=LIGHT_ACCENT, align=PP_ALIGN.CENTER, font='Calibri')

def add_image(slide, path, left, top, width=None, height=None):
    full = os.path.join(SLIDES_DIR, path)
    if not os.path.exists(full):
        add_text(slide, left, top, 4, 1, f"[Image: {path}]", size=14, color=MUTED)
        return
    kwargs = {}
    if width: kwargs['width'] = Inches(width)
    if height: kwargs['height'] = Inches(height)
    slide.shapes.add_picture(full, Inches(left), Inches(top), **kwargs)

def add_footer(slide, num):
    add_text(slide, 0.4, 7.05, 2, 0.3, "CutAndCode", size=10, color=MUTED)
    add_text(slide, 11.5, 7.05, 1.5, 0.3, f"{num} / 12", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CREAM_BOX
    shape.line.fill.background()


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image(slide, 'img-title.png', 0, 0, 13.333, 7.5)
# Title text centered
add_text(slide, 0, 2.0, 13.333, 1.0, "CutAndCode", size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font='Georgia')
add_text(slide, 0, 3.2, 13.333, 0.7, "Learning Data Structures with Your Hands", size=28, color=BROWN, align=PP_ALIGN.CENTER)
# Divider line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(4.1), Inches(1.1), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
add_text(slide, 0, 4.4, 13.333, 0.5, "Maria Fraguas  |  Creative AI  |  April 2026", size=16, color=BROWN, align=PP_ALIGN.CENTER)
add_text(slide, 11.5, 7.05, 1.5, 0.3, "1 / 12", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ==================== SLIDE 2: THE PROBLEM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_image(slide, 'img-frustrated.png', 0, 0, 13.333, 7.5)

# Title bar
title_box = add_rounded_box(slide, 0.4, 0.3, 12.5, 0.9, WARM_BG)
title_box.adjustments[0] = 0.15
tf = title_box.text_frame
tf.margin_left = Inches(0.2)
tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]
p.text = "Data Structures & Algorithms Is Hard"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = DARK
p.font.name = 'Calibri'

# Left panel
left_panel = add_rounded_box(slide, 0.4, 1.5, 6.2, 4.8, WARM_BG)
left_panel.adjustments[0] = 0.03
add_rich_bullets(slide, 0.6, 1.7, 5.8, 4.4, [
    "CS 1134 at NYU is one of the hardest intro CS courses. The material is dense and builds on itself fast.",
    "Existing resources stay abstract. Code on a screen doesn't click for everyone.",
    "Some students memorize syntax without understanding what's happening underneath.",
    'There\'s a gap between "I can follow the lecture" and "I can solve the problem."'
], size=15, spacing=10)

# Right panel
right_panel = add_rounded_box(slide, 6.9, 1.5, 5.8, 4.8, WARM_BG)
right_panel.adjustments[0] = 0.03
y = 1.8
h = add_callout_box(slide, 7.1, y, 5.4, '"The hardest part of recursion is visualizing how the stack unravels afterwards."')
y += h + 0.2
h = add_callout_box(slide, 7.1, y, 5.4, '"I can write the code but I don\'t actually understand why it works."')
y += h + 0.2
add_callout_box(slide, 7.1, y, 5.4, '"The textbook diagrams don\'t help me. I need to touch it."')

add_text(slide, 11.5, 7.05, 1.5, 0.3, "2 / 12", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ==================== SLIDE 3: THE IDEA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Physical Simulations With Household Items", size=28, bold=True)
add_rich_bullets(slide, 0.6, 1.4, 6.5, 4.5, [
    "Students pick a topic (Stacks, Sorting, Recursion...)",
    "App shows a materials list: paper, scissors, bowls, tape, pen",
    "Step-by-step tutorial guides them through physically running the algorithm",
    "For recursion: stack bowls, write sticky notes, watch the call stack build and unravel",
    "No code solutions are ever shown. Understand it physically, then code on your own."
], size=15, spacing=8)
add_image(slide, 'img-student-desk.png', 7.5, 1.1, 5.3)
add_footer(slide, 3)


# ==================== SLIDE 4: CREATIVITY & AI ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Claude as Creative Collaborator", size=28, bold=True)
add_rich_bullets(slide, 0.6, 1.3, 12, 3.6, [
    "A TA's suggestion to cut out paper to simulate algorithms sparked the idea. I started using paper strips, bowls, and household items to physically work through every algorithm in the course.",
    "I documented every learning challenge I faced and how I visualized each concept to understand it. That document became the foundation for all tutorial content.",
    "I directed the project: the ideas, the metaphors, the structure. Claude acted as my technical collaborator, building the UI and filling in implementation details from my direction.",
    "No runtime AI in the app. All content is pre-generated. Zero API keys needed."
], size=14, spacing=6)

# Flow boxes - single line
y = 5.3
bw = 2.15
flow_data = [
    ("Maria's learning", "struggles doc"),
    ("Maria designs", "physical metaphors"),
    ("Claude builds", "tutorials & UI"),
    ("Maria reviews", "& directs"),
    ("Final app", "+ content"),
]
total = 5 * bw + 4 * 0.55
sx = (13.333 - total) / 2
x = sx
for i, (l1, l2) in enumerate(flow_data):
    add_flow_box(slide, x, y, bw, 0.7, l1, l2, size=12)
    x += bw
    if i < 4:
        add_arrow(slide, x + 0.07, y + 0.08, size=24)
        x += 0.55
add_footer(slide, 4)


# ==================== SLIDE 5: TARGET USE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Who This Is For", size=28, bold=True)

add_text(slide, 0.6, 1.3, 5, 0.5, "Primary Audience", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 0.6, 1.8, 5.5, 2.0, [
    "CS 1134 students at NYU (or any intro DSA course)",
    "Visual and kinesthetic learners",
    "Students who don't learn well from lectures and textbooks alone",
], size=15, spacing=8)

add_text(slide, 0.6, 3.8, 5, 0.5, "Secondary", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 0.6, 4.3, 5.5, 0.8, [
    "TAs and tutors looking for teaching tools",
], size=15, spacing=8)

add_text(slide, 7.2, 1.3, 5, 0.5, "Accessibility First", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 7.2, 1.8, 5.5, 4.0, [
    "Free, hosted on GitHub Pages",
    "No downloads, no accounts, no frameworks",
    "Works on any device with a browser",
    "Materials: paper, scissors, tape, pen, bowls",
    "Can be used by students who don't have access to a computer to code on. An instructor could show them the instructions, they learn the algorithms first, then code on paper.",
], size=15, spacing=8)

add_footer(slide, 5)


# ==================== SLIDE 6: ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "How It's Built", size=28, bold=True)
add_rich_bullets(slide, 0.6, 1.4, 6.2, 3.5, [
    "Pure HTML, CSS, JavaScript. No frameworks.",
    "Single-page app with hash-based routing",
    "Content lives in JSON files (separate from UI)",
    "CSS is modular: variables, layout, components, tutorial styles, animations",
    "Deployable to GitHub Pages with zero config"
], size=15, spacing=8)

# Tags - moved up with more space from file tree
tags = ["HTML/CSS/JS", "ES6 Modules", "JSON Content", "GitHub Pages"]
x = 0.6
for tag in tags:
    add_tag(slide, x, 4.6, tag)
    x += 1.7

# File tree
tree_box = add_rounded_box(slide, 7.2, 1.2, 5.5, 5.0, DARK)
tree_box.adjustments[0] = 0.03
tf = tree_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.25)
tf.margin_top = Inches(0.2)
tf.margin_right = Inches(0.2)
lines = [
    ("CutAndCode/", True),
    ("  index.html", False),
    ("  css/  variables, reset, layout,", True),
    ("        components, tutorial, animations", False),
    ("  js/   app, data, ui, tutorial,", True),
    ("        components", False),
    ("  data/", True),
    ("    topics.json", False),
    ("    tutorials/", True),
    ("      stacks-operations.json", False),
    ("      stacks-midstack.json", False),
    ("      sorting-bubble.json", False),
    ("      recursion-intro.json", False),
]
for i, (line, is_dir) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x7E, 0xCB, 0xC0) if is_dir else RGBColor(0xF5, 0xB8, 0x95)
    p.font.name = 'Courier New'
    p.space_before = Pt(2)
    p.space_after = Pt(2)
add_footer(slide, 6)


# ==================== SLIDE 7: USER JOURNEY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
add_image(slide, 'img-journey.png', 0, 0, 13.333, 7.5)
add_text(slide, 11.5, 7.05, 1.5, 0.3, "7 / 12", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ==================== SLIDE 8: DOUBLY LINKED LIST ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Example: Doubly Linked List", size=28, bold=True)
add_tag(slide, 0.6, 1.2, "Core Practice", TAG_BG, TAG_FG, 1.4)
add_text(slide, 2.2, 1.22, 2, 0.3, "Coming soon", size=12, color=BROWN)

add_rich_bullets(slide, 0.6, 1.7, 5.8, 3.5, [
    "Materials: index cards, tape, two colors of string or yarn, pen",
    "Each card is a node: write data on front. Tape a blue string forward and a red string backward.",
    "Follow blue to traverse forward, red to go backward",
    "Insert by cutting both strings and re-taping through the new card",
    "Delete by connecting the strings around the removed card"
], size=14, spacing=6)

add_callout_box(slide, 0.6, 5.3, 5.5, 'You physically feel what "prev" and "next" pointers mean.')

# DLL diagram panel
dll_panel = add_rounded_box(slide, 6.6, 1.0, 6.3, 5.7, WHITE)
dll_panel.adjustments[0] = 0.03
tf = dll_panel.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_right = Inches(0.3)
tf.margin_top = Inches(0.25)

dll_lines = [
    ("two strings per card:  next \u2192  and  \u2190 prev", 14, BROWN, False),
    ("", 8, DARK, False),
    ("None  \u2190\u2192  [ 3 ]  \u2190\u2192  [ 7 ]  \u2190\u2192  [ 1 ]  \u2190\u2192  [ 9 ]  \u2190\u2192  None", 17, DARK, True),
    ("             header                                                            trailer", 11, MUTED, False),
    ("", 8, DARK, False),
    ("\u25a0 = prev (red string)          \u25a0 = next (blue string)", 12, BROWN, False),
    ("", 6, DARK, False),
    ("\u2500" * 50, 10, CREAM_BOX, False),
    ("", 6, DARK, False),
    ('inserting "5" between 7 and 1:', 15, ACCENT, False),
    ("", 6, DARK, False),
    ("[ 7 ]  \u2190\u2192  [ 5 ]  \u2190\u2192  [ 1 ]", 18, DARK, True),
    ("", 6, DARK, False),
    ("cut both strings, re-tape through the new card!", 13, MUTED, False),
]
for i, (text, size, color, bold) in enumerate(dll_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(1)
    p.space_after = Pt(1)

add_text(slide, 0.4, 7.1, 2, 0.3, "CutAndCode", size=9, color=MUTED)
add_text(slide, 11.5, 7.1, 1.5, 0.3, "8 / 12", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ==================== SLIDE 9: RECURSION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Example: The Bowl Stack Method", size=28, bold=True)
add_tag(slide, 0.6, 1.2, "Core Practice", TAG_BG, TAG_FG, 1.4)
add_text(slide, 2.2, 1.22, 3, 0.3, "15 steps, ~25 min", size=12, color=BROWN)

add_rich_bullets(slide, 0.6, 1.7, 5.8, 3.5, [
    "Materials: 4-5 bowls or cups, sticky notes, pen",
    "Each recursive call = place a new bowl on the stack",
    'Write the argument on a sticky note, mark it "WAITING"',
    "Keep stacking until you hit the base case",
    "Then unstack one at a time, filling in the return value"
], size=14, spacing=6)

add_callout_box(slide, 0.6, 5.2, 5.5, "Uses sum_to(4) as the running example")

# Placeholder for recursion diagram
add_text(slide, 7, 2.8, 5.5, 1.5, "[Insert recursion bowls diagram]", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 9)


# ==================== SLIDE 10: AI/ML ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "AI/ML in This Project", size=28, bold=True)

# Left column
add_text(slide, 0.6, 1.3, 5, 0.4, "Claude (Anthropic)", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 0.6, 1.8, 6.2, 4, [
    "Creative collaborator, not the ideas person",
    "I came up with the physical simulation metaphors. I guide Claude to create ones that make sense.",
    "Used my detailed account of learning challenges to collaborate on step-by-step tutorials. I gave the big picture, Claude fills in the details.",
    "Built the app architecture, file structure, and UI from my direction"
], size=14, spacing=8)

# Right column
add_text(slide, 7.2, 1.3, 5, 0.4, "What's NOT in the app", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 7.2, 1.8, 5.5, 2.5, [
    "No runtime AI",
    "No API keys",
    "No model calls",
    "No internet required after first load"
], size=14, spacing=8)

add_callout_box(slide, 7.2, 4.5, 5.5, 'This is a "Creative AI" project in the truest sense: AI as a design partner, not a feature.')
add_footer(slide, 10)


# ==================== SLIDE 11: PROTOTYPE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Current State of the Prototype", size=28, bold=True)

add_text(slide, 0.6, 1.3, 5, 0.4, "What's Working", size=18, bold=True, color=BROWN)
add_rich_bullets(slide, 0.6, 1.8, 5.5, 4, [
    "4 tutorials complete across 3 topics",
    "Landing page, topic grid, tutorial list",
    "Materials checklist screen",
    "Step-by-step navigation",
    "CSS-rendered diagrams",
    "Collapsible concept notes, hints",
    "Reflection flashcards",
    "Responsive layout, hash routing"
], size=13, spacing=5)

# Tags
tags = ["3 Topics", "4 Tutorials", "21 Topics Mapped"]
x = 0.6
for tag in tags:
    add_tag(slide, x, 5.6, tag)
    x += 1.85

# Right side - placeholder
add_text(slide, 7, 3.2, 5.5, 1, "[Insert screenshots here]", size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 11)


# ==================== SLIDE 12: CHALLENGES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_gradient_bg(slide)
add_text(slide, 0.6, 0.4, 10, 0.8, "Challenges & Questions", size=28, bold=True)

add_text(slide, 0.6, 1.3, 5, 0.4, "Challenges", size=18, bold=True, color=RED_ACCENT)
add_rich_bullets(slide, 0.6, 1.8, 6, 5, [
    "Scaling content: 4 tutorials done, 21 topics to cover. Each needs a unique physical metaphor.",
    "Primary challenge: making sure the doodle-style drawings work. I need to review each drawing for every exercise to make sure text isn't covered by other elements and everything is human-readable.",
    "I need to physically simulate every algorithm myself to verify the steps work before publishing.",
    "Making visuals appealing enough that it doesn't look like every other algorithm tutorial with moving boxes and numbers on a screen."
], size=13, spacing=6)

add_text(slide, 7.2, 1.3, 5, 0.4, "Questions for You", size=18, bold=True, color=BLUE_ACCENT)
add_rich_bullets(slide, 7.2, 1.8, 5.5, 5, [
    "How does the app feel to click around on?",
    "For topics like Big-O and amortized analysis: do you have ideas for how to make those physical, or should I leave them out since they're not algorithm-specific?",
    "How far should the doodle-style visuals go? Full illustrated guides per step, or just key moments?",
    "Any recommendations for creating hand-drawn SVG illustrations efficiently?"
], size=13, spacing=6)
add_footer(slide, 12)


# Save
output_path = os.path.join(SLIDES_DIR, 'CutAndCode_ProjectPlan.pptx')
prs.save(output_path)
print(f"Saved to {output_path}")
